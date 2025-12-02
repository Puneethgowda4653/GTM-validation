"""
GTM Validation Engine - Enhanced with Supabase
Beautiful UI/UX with Flask, Supabase backend, and historical tracking
"""

from flask import Flask, render_template, request, jsonify, send_file
import os
import json
from pathlib import Path
import google.generativeai as genai
import PyPDF2
from pptx import Presentation
from werkzeug.utils import secure_filename
import tempfile
from datetime import datetime
from dotenv import load_dotenv
from supabase import create_client, Client

# Load environment variables
load_dotenv()

app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = int(os.getenv('MAX_CONTENT_LENGTH', 52428800))
app.config['UPLOAD_FOLDER'] = os.getenv('UPLOAD_FOLDER', tempfile.gettempdir())
app.config['SECRET_KEY'] = os.getenv('SECRET_KEY', 'dev-secret-key-change-in-production')

# Create upload folder if it doesn't exist
if app.config['UPLOAD_FOLDER'] != tempfile.gettempdir():
    os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

# Initialize Supabase
supabase_url = os.getenv('SUPABASE_URL')
supabase_key = os.getenv('SUPABASE_KEY')
supabase: Client = None

if supabase_url and supabase_key:
    supabase = create_client(supabase_url, supabase_key)
    print("✅ Supabase connected")
else:
    print("⚠️  Supabase not configured - running without database")

class GTMValidator:
    def __init__(self, api_key: str, model_name: str = None):
        genai.configure(api_key=api_key)
        model_name = model_name or os.getenv('GEMINI_MODEL', 'gemini-2.5-flash')
        self.model = genai.GenerativeModel(model_name)
        print(f"✅ Using model: {model_name}")
        
        # Updated framework with 1-5 scoring
        self.framework = {
            "Business Clarity": {
                "weight": 10,
                "criteria": [
                    "Defined target customer",
                    "Clear value proposition",
                    "Aligned problem statement"
                ]
            },
            "GTM Hypothesis": {
                "weight": 15,
                "criteria": [
                    "Defined customer segment",
                    "Channel strategy stated",
                    "Pricing model defined",
                    "Buyer persona identified"
                ]
            },
            "Market Validation": {
                "weight": 15,
                "criteria": [
                    "User interviews done",
                    "Evidence of problem urgency",
                    "Signups / trials",
                    "Customer interest rate >30%"
                ]
            },
            "Product Validation": {
                "weight": 15,
                "criteria": [
                    "Demo → Customer conversion ≥15%",
                    "Retention ≥50%"
                ]
            },
            "Channel & Messaging": {
                "weight": 15,
                "criteria": [
                    "Top channel identified",
                    "Click-Through Rate >2–5%",
                    "Lead qualification >40%",
                    "Sales cycle defined"
                ]
            },
            "Financial Validation": {
                "weight": 10,
                "criteria": [
                    "CAC declining",
                    "LTV ≥ 3× CAC",
                    "Payback < 12 months",
                    "Margin sustainable"
                ]
            },
            "Operational Readiness": {
                "weight": 10,
                "criteria": [
                    "Sales playbook ready",
                    "Delivery team ready",
                    "Onboarding defined",
                    "Support processes in place"
                ]
            },
            "Experimentation": {
                "weight": 10,
                "criteria": [
                    "3+ GTM experiments run",
                    "Results documented",
                    "Learnings implemented"
                ]
            }
        }

    def extract_text_from_pdf(self, file_path: str) -> str:
        text = ""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
        except Exception as e:
            raise Exception(f"Error extracting PDF: {e}")
        return text

    def extract_text_from_pptx(self, file_path: str) -> str:
        text = ""
        try:
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            raise Exception(f"Error extracting PPTX: {e}")
        return text

    def extract_content(self, file_path: str) -> str:
        file_extension = Path(file_path).suffix.lower()
        
        if file_extension == '.pdf':
            return self.extract_text_from_pdf(file_path)
        elif file_extension in ['.pptx', '.ppt']:
            return self.extract_text_from_pptx(file_path)
        elif file_extension == '.txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        else:
            raise ValueError(f"Unsupported file format: {file_extension}")

    def create_validation_prompt(self, deck_content: str, company_name: str) -> str:
        prompt = f"""You are an advanced GTM Validation Engine. Analyze this pitch deck for {company_name}.

PITCH DECK CONTENT:
{deck_content}

SCORING SYSTEM (1-5):
1 = No clarity / no validation
2 = Idea-level hypothesis
3 = Early validation or pilot stage
4 = Clear traction / partial validation
5 = Fully validated & scalable

GTM DIMENSIONS & WEIGHTS:
{json.dumps(self.framework, indent=2)}

INSTRUCTIONS:
1. For EACH dimension, provide a score from 1-5 based on evidence
2. Provide TWO types of summaries:
   - key_summary: 10-12 words max, concise insight (e.g., "Strong clarity from founder + mentors agree on product identity")
   - detailed_summary: Full explanation with evidence from deck
3. Calculate weighted scores: score × (weight/100)
   Example: score=3, weight=10% → weighted_score = 3 × 0.10 = 0.30
4. Determine overall stage/readiness

OUTPUT FORMAT (MUST BE VALID JSON):
{{
  "company_name": "{company_name}",
  "dimensions": [
    {{
      "name": "Business Clarity",
      "weight": 10,
      "score": 3,
      "key_summary": "Strong clarity from founder + mentors agree on product identity",
      "detailed_summary": "Detailed explanation of why this score, citing specific evidence from deck",
      "weighted_score": 0.30,
      "criteria_met": ["criterion 1", "criterion 2"]
    }}
  ],
  "total_weighted_score": 0.0,
  "total_score_out_of_5": 0.0,
  "stage_readiness": "Iterate & Refine",
  "overall_summary": "Comprehensive assessment of GTM readiness"
}}

WEIGHTED SCORE CALCULATION:
- Formula: score × (weight ÷ 100)
- Business Clarity (10%): score × 0.10
- GTM Hypothesis (15%): score × 0.15
- Market Validation (15%): score × 0.15
- Product Validation (15%): score × 0.15
- Channel & Messaging (15%): score × 0.15
- Financial Validation (10%): score × 0.10
- Operational Readiness (10%): score × 0.10
- Experimentation (10%): score × 0.10

STAGE/READINESS MAPPING:
- Score 4-5: "Validated & Scalable" 
- Score 3-3.9: "Iterate & Refine"
- Score 2-2.9: "Early Stage"
- Score 0-1: "Pivot Required"

Respond ONLY with valid JSON."""
        
        return prompt

    def validate_deck(self, file_path: str, company_name: str) -> dict:
        deck_content = self.extract_content(file_path)
        
        if not deck_content.strip():
            return {"error": "No content extracted from the file"}
        
        prompt = self.create_validation_prompt(deck_content, company_name)
        
        try:
            response = self.model.generate_content(prompt)
            result_text = response.text.strip()
            
            if result_text.startswith('```'):
                result_text = result_text.split('```')[1]
                if result_text.startswith('json'):
                    result_text = result_text[4:]
                result_text = result_text.strip()
            
            result = json.loads(result_text)
            result['validated_date'] = datetime.now().isoformat()
            
            return result
            
        except json.JSONDecodeError as e:
            return {"error": "Failed to parse AI response", "raw": response.text[:500]}
        except Exception as e:
            return {"error": str(e)}

def save_to_supabase(validation_result: dict) -> dict:
    """Save validation result to Supabase"""
    if not supabase:
        print("❌ ERROR: Supabase not configured!")
        print("Check your .env file for SUPABASE_URL and SUPABASE_KEY")
        return {"error": "Supabase not configured"}
    
    try:
        company_name = validation_result.get('company_name', 'Unknown Company')
        print(f"\n{'='*60}")
        print(f"💾 SAVING TO SUPABASE: {company_name}")
        print(f"{'='*60}")
        
        # Prepare validation data
        validation_data = {
            "company_name": company_name,
            "total_weighted_score": float(validation_result.get('total_weighted_score', 0)),
            "total_score_out_of_5": float(validation_result.get('total_score_out_of_5', 0)),
            "stage_readiness": str(validation_result.get('stage_readiness', '')),
            "overall_summary": str(validation_result.get('overall_summary', '')),
            "validated_date": validation_result.get('validated_date')
        }
        
        print(f"📊 Base data prepared:")
        print(f"   - Total Score: {validation_data['total_score_out_of_5']}")
        print(f"   - Stage: {validation_data['stage_readiness']}")
        
        # Add dimension scores
        dimensions_added = 0
        for dim in validation_result.get('dimensions', []):
            dim_name = dim['name'].lower().replace(' ', '_').replace('&', 'and').replace('  ', '_')
            score = int(dim.get('score', 0))
            weighted = float(dim.get('weighted_score', 0))
            summary = str(dim.get('key_summary', '') or dim.get('summary', ''))
            
            validation_data[f"{dim_name}_score"] = score
            validation_data[f"{dim_name}_summary"] = summary
            validation_data[f"{dim_name}_weighted"] = weighted
            
            dimensions_added += 1
            print(f"   - {dim['name']}: Score={score}, Weighted={weighted:.2f}")
        
        print(f"✅ Added {dimensions_added} dimensions")
        print(f"📤 Inserting into Supabase...")
        
        # Insert validation (skip companies table for now)
        validation_response = supabase.table('validations').insert(validation_data).execute()
        
        if validation_response.data and len(validation_response.data) > 0:
            record_id = validation_response.data[0].get('id', 'N/A')
            print(f"✅ SUCCESS! Saved with ID: {record_id}")
            print(f"{'='*60}\n")
            return {"success": True, "id": record_id}
        else:
            print(f"⚠️  WARNING: Insert returned no data")
            print(f"Response: {validation_response}")
            return {"error": "Insert returned no data"}
        
    except Exception as e:
        print(f"\n{'='*60}")
        print(f"❌ ERROR SAVING TO SUPABASE!")
        print(f"{'='*60}")
        print(f"Error: {str(e)}")
        print(f"Error Type: {type(e).__name__}")
        import traceback
        traceback.print_exc()
        print(f"{'='*60}\n")
        return {"error": f"Database error: {str(e)}"}

def get_validation_history():
    """Get all validation history from Supabase"""
    if not supabase:
        print("⚠️  Supabase not configured")
        return []
    
    try:
        response = supabase.table('validations').select('*').order('validated_date', desc=True).execute()
        print(f"✅ Database query successful: {len(response.data)} records found")
        
        # Debug: print first record if exists
        if response.data and len(response.data) > 0:
            print(f"📊 Sample record: {response.data[0].get('company_name', 'N/A')}")
        
        return response.data
    except Exception as e:
        print(f"❌ Error fetching history: {e}")
        import traceback
        traceback.print_exc()
        return []

# Global validator instance
validator = None

@app.route('/')
def index():
    has_env_key = bool(os.getenv('GEMINI_API_KEY'))
    has_supabase = bool(supabase)
    return render_template('index.html', has_env_key=has_env_key, has_supabase=has_supabase)

@app.route('/api/validate', methods=['POST'])
def validate():
    global validator
    
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400
    
    file = request.files['file']
    company_name = request.form.get('company_name', 'Unknown Company')
    
    # Get API key from form or environment
    api_key = request.form.get('api_key') or os.getenv('GEMINI_API_KEY')
    
    if not api_key:
        return jsonify({'error': 'API key is required'}), 400
    
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
    
    # Save file temporarily
    filename = secure_filename(file.filename)
    filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(filepath)
    
    try:
        # Initialize validator with API key
        validator = GTMValidator(api_key)
        
        # Validate
        result = validator.validate_deck(filepath, company_name)
        
        # Clean up file
        os.remove(filepath)
        
        if 'error' not in result:
            # Save to Supabase
            db_result = save_to_supabase(result)
            result['database_saved'] = db_result.get('success', False)
        
        return jsonify(result)
    
    except Exception as e:
        if os.path.exists(filepath):
            os.remove(filepath)
        return jsonify({'error': str(e)}), 500

@app.route('/api/history', methods=['GET'])
def get_history():
    """Get validation history"""
    if not supabase:
        return jsonify({"error": "Database not configured"}), 500
    
    try:
        history = get_validation_history()
        print(f"✅ Fetched {len(history)} validation records")
        return jsonify(history)
    except Exception as e:
        print(f"❌ Error fetching history: {e}")
        return jsonify({"error": str(e)}), 500

@app.route('/api/download-report', methods=['POST'])
def download_report():
    data = request.json
    
    company_name = data.get('company_name', 'Company')
    
    # Generate text report
    report_lines = []
    report_lines.append("=" * 80)
    report_lines.append(f"GTM SCOREBOARD — {company_name}")
    report_lines.append(f"Validated: {data.get('validated_date', '')}")
    report_lines.append("=" * 80)
    report_lines.append("")
    
    # Table header
    report_lines.append(f"{'GTM Criteria':<30} {'Weight':<10} {'Score':<10} {'Weighted Score':<15}")
    report_lines.append("-" * 80)
    
    for dim in data.get('dimensions', []):
        key_summary = dim.get('key_summary', dim.get('summary', ''))
        detailed_summary = dim.get('detailed_summary', dim.get('summary', ''))
        weighted = dim['weighted_score']
        
        report_lines.append(f"{dim['name']:<30} {dim['weight']}%{'':<7} {dim['score']}/5{'':<5} {weighted:<15.2f}")
        report_lines.append(f"  Key: {key_summary}")
        if detailed_summary and detailed_summary != key_summary:
            report_lines.append(f"  Details: {detailed_summary}")
        report_lines.append("")
    
    report_lines.append("=" * 80)
    report_lines.append(f"Total Weighted GTM Score: {data['total_weighted_score']:.2f} / 100")
    report_lines.append(f"Overall Score: {data['total_score_out_of_5']:.2f} / 5")
    report_lines.append(f"Stage/Readiness: {data['stage_readiness']}")
    report_lines.append("=" * 80)
    report_lines.append(f"\nOVERALL SUMMARY:\n{data['overall_summary']}")
    
    report = "\n".join(report_lines)
    
    # Save to temp file
    temp_file = os.path.join(tempfile.gettempdir(), f'gtm_report_{company_name}.txt')
    with open(temp_file, 'w', encoding='utf-8') as f:
        f.write(report)
    
    return send_file(temp_file, as_attachment=True, download_name=f'gtm_report_{company_name}.txt')

if __name__ == '__main__':
    print("=" * 80)
    print("🚀 GTM VALIDATION ENGINE - ENHANCED")
    print("=" * 80)
    
    # Check for API key
    api_key = os.getenv('GEMINI_API_KEY')
    if api_key:
        print("✅ API Key loaded from .env file")
    else:
        print("⚠️  No API Key in .env - users will need to enter manually")
    
    # Check Supabase
    if supabase:
        print("✅ Supabase connected - database enabled")
    else:
        print("⚠️  Supabase not configured - running without database")
    
    host = os.getenv('HOST', '0.0.0.0')
    port = int(os.getenv('PORT', 5000))
    
    print(f"📱 Open your browser and go to: http://localhost:{port}")
    print("=" * 80)
    
    app.run(debug=True, host=host, port=port)